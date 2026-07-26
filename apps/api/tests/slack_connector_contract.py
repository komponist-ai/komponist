"""Provider-free contract checks for the organization-scoped Slack connector."""

import asyncio
import io
import json
import sys
from datetime import datetime
from pathlib import Path
from unittest.mock import patch
from urllib.parse import parse_qs, urlsplit

from fastapi import HTTPException
from starlette.requests import Request

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
sys.path.insert(0, str(Path(__file__).resolve().parents[3] / "packages"))

from integrations.slack import (  # noqa: E402
    _extract_slack_file_text,
    fetch_slack_threads,
    get_oauth_url,
    handle_slack_webhook,
    list_channels,
)
from core.models import SourceItem, SourceType  # noqa: E402


class FakeResponse:
    status_code = 200

    def __init__(self, payload: dict | None = None, *, content: bytes = b""):
        self._payload = payload
        self.content = content

    def json(self) -> dict:
        assert self._payload is not None
        return self._payload


class FakeSlackClient:
    seen_authorizations: list[str] = []

    def __init__(self, *, headers: dict, **_kwargs):
        self.headers = headers
        self.seen_authorizations.append(headers["Authorization"])

    async def __aenter__(self):
        return self

    async def __aexit__(self, *_args):
        return None

    async def get(self, url: str, params: dict | None = None):
        method = url.rsplit("/", 1)[-1]
        if method == "strategy.md":
            return FakeResponse(
                content=(
                    b"# Campus strategy\n\n"
                    b"Decision: The board approved a four-week member campaign."
                ),
            )
        if method == "conversations.list":
            return FakeResponse({
                "ok": True,
                "channels": [{
                    "id": "C123",
                    "name": "board",
                    "is_private": True,
                    "is_member": True,
                }],
                "response_metadata": {"next_cursor": ""},
            })
        if method == "conversations.history":
            assert params["channel"] == "C123"
            return FakeResponse({
                "ok": True,
                "messages": [{"ts": "1784970000.000100"}],
                "response_metadata": {"next_cursor": ""},
            })
        if method == "conversations.replies":
            return FakeResponse({
                "ok": True,
                "messages": [
                    {
                        "ts": "1784970000.000100",
                        "user": "U1",
                        "text": "The board approved the event budget.",
                        "files": [{
                            "id": "F123",
                            "name": "strategy.md",
                            "size": 78,
                            "user": "U1",
                            "timestamp": 1784970010,
                            "url_private_download": "https://files.slack.test/strategy.md",
                            "permalink": "https://workspace.slack.com/files/U1/F123",
                        }],
                    },
                    {
                        "ts": "1784970030.000200",
                        "user": "U2",
                        "text": "The limit is EUR 4,800.",
                    },
                ],
            })
        if method == "users.info":
            name = "Amir" if params["user"] == "U1" else "Lea"
            return FakeResponse({
                "ok": True,
                "user": {"real_name": name, "profile": {}},
            })
        raise AssertionError(f"Unexpected Slack method: {method}")


async def check_connector() -> None:
    with patch("integrations.slack.httpx.AsyncClient", FakeSlackClient):
        channels = await list_channels("xoxb-org-token")
        assert channels == [{
            "id": "C123",
            "name": "board",
            "is_private": True,
            "is_member": True,
        }]
        items = await fetch_slack_threads(
            "org-campus",
            access_token="xoxb-org-token",
            channel_ids=["C123"],
            channel_names={"C123": "board"},
            department_id="department-board",
        )

    assert len(items) == 2
    by_reference = {item.reference: item for item in items}
    thread = by_reference["slack:C123/1784970000.000100"]
    attachment = by_reference["slack-file:F123"]
    assert thread.title.startswith("#board —")
    assert "Amir: The board approved" in thread.body
    assert "[Attached: strategy.md]" in thread.body
    assert "Lea: The limit is EUR 4,800." in thread.body
    assert attachment.title == "strategy.md"
    assert "four-week member campaign" in attachment.body
    assert attachment.kind == "file"
    assert attachment.url == "https://workspace.slack.com/files/U1/F123"
    assert all(item.department_id == "department-board" for item in items)
    assert set(FakeSlackClient.seen_authorizations) == {"Bearer xoxb-org-token"}

    query = parse_qs(urlsplit(get_oauth_url("safe-state")).query)
    scopes = set(query["scope"][0].split(","))
    assert {"channels:history", "groups:history", "files:read", "users:read"} <= scopes
    assert query["state"] == ["safe-state"]
    print("✓ Slack OAuth, channel discovery, thread and attachment assembly, and org token scope")


async def check_sync_handoff() -> None:
    import main

    captured: dict = {}

    async def fake_fetch(org_id: str, **kwargs):
        captured.update({"org_id": org_id, **kwargs})
        return [SourceItem(
            org_id=org_id,
            source=SourceType.SLACK,
            kind="thread",
            title="#board — Budget approval",
            body="Amir: The board approved EUR 4,800.",
            url="https://slack.com/archives/C123/p1784970000000100",
            reference="slack:C123/1784970000.000100",
            source_date=datetime(2026, 7, 25),
            department_id=kwargs["department_id"],
        )]

    async def fake_settings(_org_id: str):
        return {"auto_confirm": False, "parallel_batch_size": 5}

    async def fake_extraction(source_item, auto_confirm: bool = False):
        assert source_item.department_id == "department-board"
        assert auto_confirm is False
        return {"entities_created": 2, "relationships_created": 1}

    with (
        patch("integrations.slack.fetch_slack_threads", fake_fetch),
        patch.object(main, "get_org_settings", fake_settings),
        patch.object(main, "run_extraction", fake_extraction),
    ):
        result = await main.sync_slack_source(
            "org-campus",
            {
                "departmentId": "department-board",
                "config": {
                    "token": "xoxb-org-token",
                    "watched_channels": ["C123"],
                    "channel_names": {"C123": "board"},
                },
            },
        )

    assert captured["access_token"] == "xoxb-org-token"
    assert captured["channel_ids"] == ["C123"]
    assert captured["department_id"] == "department-board"
    assert result["items_processed"] == 1
    assert result["entities_created"] == 2
    assert result["relationships_created"] == 1
    print("✓ manual sync hands selected Slack threads to extraction and review")


def check_attachment_parsers() -> None:
    from docx import Document
    from pptx import Presentation
    from reportlab.pdfgen import canvas

    docx_buffer = io.BytesIO()
    docx_document = Document()
    docx_document.add_paragraph("The finance department approved EUR 4,800.")
    docx_document.save(docx_buffer)
    assert "approved EUR 4,800" in _extract_slack_file_text(
        docx_buffer.getvalue(), "budget.docx"
    )

    pptx_buffer = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(0, 0, 3_000_000, 800_000)
    textbox.text = "Campus campaign runs for four weeks."
    presentation.save(pptx_buffer)
    assert "four weeks" in _extract_slack_file_text(
        pptx_buffer.getvalue(), "campaign.pptx"
    )

    pdf_buffer = io.BytesIO()
    pdf = canvas.Canvas(pdf_buffer)
    pdf.drawString(72, 720, "The board confirmed the sponsorship policy.")
    pdf.save()
    assert "sponsorship policy" in _extract_slack_file_text(
        pdf_buffer.getvalue(), "policy.pdf"
    )
    print("✓ PDF, DOCX, and PPTX Slack attachment parsers return extractable text")


def slack_request(payload: dict) -> Request:
    body = json.dumps(payload).encode("utf-8")
    delivered = False

    async def receive():
        nonlocal delivered
        if delivered:
            return {"type": "http.disconnect"}
        delivered = True
        return {"type": "http.request", "body": body, "more_body": False}

    return Request({
        "type": "http",
        "method": "POST",
        "path": "/webhooks/slack",
        "headers": [],
        "query_string": b"",
    }, receive)


async def check_webhook_tenant_binding() -> None:
    async def fake_sources(_org_id: str, include_config: bool = False):
        assert include_config is True
        return [{
            "type": "slack",
            "config": {"team_id": "T-CAMPUS"},
        }]

    with (
        patch("integrations.slack.verify_slack_signature", return_value=True),
        patch("persistence.list_connected_sources", fake_sources),
    ):
        challenge = await handle_slack_webhook(
            slack_request({
                "type": "url_verification",
                "team_id": "T-CAMPUS",
                "challenge": "verified",
            }),
            "org-campus",
        )
        assert challenge == {"challenge": "verified"}

        try:
            await handle_slack_webhook(
                slack_request({
                    "type": "url_verification",
                    "team_id": "T-OTHER",
                    "challenge": "wrong-org",
                }),
                "org-campus",
            )
        except HTTPException as error:
            assert error.status_code == 403
        else:
            raise AssertionError("A signed event from another workspace was accepted")
    print("✓ signed Slack events remain bound to their Komponist organization")


if __name__ == "__main__":
    asyncio.run(check_connector())
    check_attachment_parsers()
    asyncio.run(check_sync_handoff())
    asyncio.run(check_webhook_tenant_binding())
    print("Slack connector contract: OK")
