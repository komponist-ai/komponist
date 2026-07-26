"""
Slack integration.

Webhook handling, thread assembly, backfill, and normalization to SourceItem.
"""

import hmac
import hashlib
import io
import os
import json
from datetime import datetime, timedelta
from typing import Optional, Dict, Any, List
from collections import defaultdict

import httpx
from fastapi import Request, HTTPException

import sys
sys.path.append("../../../packages")

from core.models import SourceItem, SourceType
from database import async_session, EventRaw, SyncState


SLACK_SIGNING_SECRET = os.getenv("SLACK_SIGNING_SECRET", "")
SLACK_BOT_TOKEN = os.getenv("SLACK_BOT_TOKEN", "")
SLACK_CLIENT_ID = os.getenv("SLACK_CLIENT_ID", "")
SLACK_CLIENT_SECRET = os.getenv("SLACK_CLIENT_SECRET", "")
SLACK_REDIRECT_URI = os.getenv("SLACK_REDIRECT_URI", "http://localhost:8000/auth/slack/callback")
MAX_SLACK_FILE_BYTES = 10 * 1024 * 1024
SLACK_FILE_EXTENSIONS = {
    ".csv", ".docx", ".json", ".markdown", ".md", ".pdf", ".pptx",
    ".tsv", ".txt", ".yaml", ".yml",
}

# Watched channels (configured per org)
WATCHED_CHANNELS = [
    channel.strip()
    for channel in os.getenv("SLACK_WATCHED_CHANNELS", "").split(",")
    if channel.strip()
]


class SlackApiError(RuntimeError):
    """A safe, user-facing Slack API failure."""


def _headers(access_token: str) -> dict[str, str]:
    token = access_token.strip()
    if not token:
        raise SlackApiError("Slack access token is missing")
    return {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }


async def _get(
    client: httpx.AsyncClient,
    method: str,
    params: dict[str, Any],
) -> dict[str, Any]:
    response = await client.get(f"https://slack.com/api/{method}", params=params)
    try:
        payload = response.json()
    except ValueError as error:
        raise SlackApiError("Slack returned an unreadable response") from error
    if response.status_code != 200 or not payload.get("ok"):
        reason = str(payload.get("error") or f"HTTP {response.status_code}")
        raise SlackApiError(f"Slack API error: {reason}")
    return payload


def get_oauth_url(state: str) -> str:
    """
    Generate Slack OAuth URL.

    Args:
        state: State parameter for CSRF protection (should include org_id)

    Returns:
        OAuth authorization URL
    """
    scopes = ",".join([
        "channels:history",
        "channels:read",
        "groups:history",
        "groups:read",
        "files:read",
        "users:read",
    ])
    params = {
        "client_id": SLACK_CLIENT_ID,
        "redirect_uri": SLACK_REDIRECT_URI,
        "scope": scopes,
        "state": state,
    }
    from urllib.parse import urlencode
    query = urlencode(params)
    return f"https://slack.com/oauth/v2/authorize?{query}"


async def list_channels(access_token: str) -> List[Dict[str, Any]]:
    """Return channels the installed bot can currently read."""
    channels: list[dict[str, Any]] = []
    cursor: Optional[str] = None
    async with httpx.AsyncClient(headers=_headers(access_token), timeout=30) as client:
        while True:
            params: dict[str, Any] = {
                "types": "public_channel,private_channel",
                "exclude_archived": "true",
                "limit": 200,
            }
            if cursor:
                params["cursor"] = cursor
            payload = await _get(client, "conversations.list", params)
            channels.extend(
                {
                    "id": channel["id"],
                    "name": channel.get("name") or channel["id"],
                    "is_private": bool(channel.get("is_private")),
                    "is_member": bool(channel.get("is_member")),
                }
                for channel in payload.get("channels", [])
                if channel.get("id")
            )
            cursor = payload.get("response_metadata", {}).get("next_cursor") or None
            if not cursor:
                break
    return sorted(channels, key=lambda channel: channel["name"].casefold())


async def exchange_code(code: str) -> Dict[str, Any]:
    """
    Exchange authorization code for access token.

    Args:
        code: Authorization code from OAuth callback

    Returns:
        Token response with access_token, team info, etc.
    """
    async with httpx.AsyncClient() as client:
        response = await client.post(
            "https://slack.com/api/oauth.v2.access",
            data={
                "client_id": SLACK_CLIENT_ID,
                "client_secret": SLACK_CLIENT_SECRET,
                "code": code,
                "redirect_uri": SLACK_REDIRECT_URI,
            }
        )

        data = response.json()

        if not data.get("ok"):
            raise HTTPException(status_code=400, detail=data.get("error", "OAuth failed"))

        return data


def verify_slack_signature(body: bytes, timestamp: str, signature: str) -> bool:
    """
    Verify Slack request signature.

    Args:
        body: Raw request body
        timestamp: X-Slack-Request-Timestamp header
        signature: X-Slack-Signature header

    Returns:
        True if signature is valid
    """
    if not SLACK_SIGNING_SECRET:
        return os.getenv("KOMPONIST_ALLOW_UNSIGNED_WEBHOOKS", "false").lower() == "true"

    # Check timestamp (prevent replay attacks)
    try:
        request_timestamp = int(timestamp)
    except (TypeError, ValueError):
        return False
    current_timestamp = int(datetime.utcnow().timestamp())
    if abs(current_timestamp - request_timestamp) > 60 * 5:
        return False

    sig_basestring = f"v0:{timestamp}:{body.decode()}"
    expected_signature = "v0=" + hmac.new(
        SLACK_SIGNING_SECRET.encode(),
        sig_basestring.encode(),
        hashlib.sha256
    ).hexdigest()

    return hmac.compare_digest(expected_signature, signature)


async def handle_slack_webhook(request: Request, org_id: str) -> Dict[str, Any]:
    """
    Handle Slack webhook event.

    Args:
        request: FastAPI request
        org_id: Organization ID

    Returns:
        Status dict or challenge response
    """
    timestamp = request.headers.get("X-Slack-Request-Timestamp", "")
    signature = request.headers.get("X-Slack-Signature", "")

    body = await request.body()

    if not verify_slack_signature(body, timestamp, signature):
        raise HTTPException(status_code=401, detail="Invalid signature")

    try:
        payload = json.loads(body)
    except (UnicodeDecodeError, json.JSONDecodeError) as error:
        raise HTTPException(status_code=400, detail="Invalid JSON payload") from error

    # A valid signature proves the request came from the configured Slack app,
    # but not which Komponist organization owns the workspace. Bind the signed
    # team id to the encrypted connector config before accepting any event.
    from persistence import list_connected_sources

    sources = await list_connected_sources(org_id, include_config=True)
    source = next((item for item in sources if item["type"] == "slack"), None)
    configured_team_id = (
        source.get("config", {}).get("team_id")
        if source
        else None
    )
    if (
        not configured_team_id
        or payload.get("team_id") != configured_team_id
    ):
        raise HTTPException(
            status_code=403,
            detail="Slack workspace does not match this organization",
        )

    # Handle URL verification challenge
    if payload.get("type") == "url_verification":
        return {"challenge": payload.get("challenge")}

    # Store raw event
    event = payload.get("event", {})
    event_type = event.get("type", "unknown")

    async with async_session() as session:
        raw_event = EventRaw(
            org_id=org_id,
            source="slack",
            event_type=event_type,
            payload=payload
        )
        session.add(raw_event)
        await session.commit()

    return {"status": "received", "event_type": event_type}


async def assemble_thread(
    channel: str,
    thread_ts: str,
    org_id: str,
    *,
    access_token: Optional[str] = None,
    channel_name: Optional[str] = None,
    department_id: Optional[str] = None,
) -> Optional[SourceItem]:
    """
    Assemble a complete Slack thread into a SourceItem.

    Args:
        channel: Channel ID
        thread_ts: Thread timestamp
        org_id: Organization ID

    Returns:
        SourceItem with full thread content
    """
    items = await assemble_thread_items(
        channel,
        thread_ts,
        org_id,
        access_token=access_token,
        channel_name=channel_name,
        department_id=department_id,
    )
    return items[0] if items else None


def _extract_slack_file_text(content: bytes, filename: str) -> str:
    """Extract text from the document formats accepted by the Slack connector."""
    suffix = os.path.splitext(filename)[1].lower()
    if suffix not in SLACK_FILE_EXTENSIONS:
        raise SlackApiError(
            f"Unsupported Slack file type for {filename}. "
            "Supported: Markdown, text, YAML, JSON, CSV, PDF, DOCX, PPTX"
        )

    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
    elif suffix == ".docx":
        from docx import Document

        document = Document(io.BytesIO(content))
        text = "\n\n".join(
            paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()
        )
    elif suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        slides = []
        for slide in presentation.slides:
            parts = [
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and str(shape.text).strip()
            ]
            if parts:
                slides.append("\n".join(parts))
        text = "\n\n".join(slides)
    else:
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError as error:
            raise SlackApiError(f"{filename} is not UTF-8 encoded text") from error

    if not text.strip():
        raise SlackApiError(
            f"{filename} contains no extractable text. Scanned PDFs need OCR first."
        )
    return text


async def _download_slack_file(
    access_token: str,
    file_data: Dict[str, Any],
) -> tuple[str, bytes]:
    """Download one private Slack file with the workspace OAuth token."""
    filename = str(file_data.get("name") or f"{file_data.get('id', 'file')}.txt")
    suffix = os.path.splitext(filename)[1].lower()
    if suffix not in SLACK_FILE_EXTENSIONS:
        raise SlackApiError(
            f"Unsupported Slack file type for {filename}. "
            "Supported: Markdown, text, YAML, JSON, CSV, PDF, DOCX, PPTX"
        )
    declared_size = int(file_data.get("size") or 0)
    if declared_size > MAX_SLACK_FILE_BYTES:
        raise SlackApiError(f"{filename} exceeds the 10 MB Slack file limit")

    download_url = (
        file_data.get("url_private_download")
        or file_data.get("url_private")
    )
    if not download_url:
        raise SlackApiError(f"Slack did not provide a download URL for {filename}")

    async with httpx.AsyncClient(
        headers={"Authorization": f"Bearer {access_token}"},
        timeout=60,
        follow_redirects=True,
    ) as client:
        response = await client.get(str(download_url))
    if response.status_code != 200:
        raise SlackApiError(
            f"Slack could not download {filename} (HTTP {response.status_code})"
        )
    content = response.content
    if len(content) > MAX_SLACK_FILE_BYTES:
        raise SlackApiError(f"{filename} exceeds the 10 MB Slack file limit")
    return filename, content


async def assemble_thread_items(
    channel: str,
    thread_ts: str,
    org_id: str,
    *,
    access_token: Optional[str] = None,
    channel_name: Optional[str] = None,
    department_id: Optional[str] = None,
) -> List[SourceItem]:
    """Assemble a Slack thread plus each supported attached document."""
    token = (access_token or SLACK_BOT_TOKEN).strip()
    if not token:
        return []

    async with httpx.AsyncClient(headers=_headers(token), timeout=30) as client:
        # Fetch thread messages
        data = await _get(
            client,
            "conversations.replies",
            {"channel": channel, "ts": thread_ts, "limit": 200},
        )

        messages = data.get("messages", [])

        if not messages:
            return []

        # Resolve user names
        user_cache = {}

        async def get_user_name(user_id: str) -> str:
            if user_id in user_cache:
                return user_cache[user_id]

            try:
                user_data = await _get(client, "users.info", {"user": user_id})
            except SlackApiError:
                return user_id
            name = (
                user_data.get("user", {}).get("profile", {}).get("display_name")
                or user_data.get("user", {}).get("real_name")
                or user_id
            )
            user_cache[user_id] = name
            return name

        # Build thread body
        thread_lines = []

        for msg in messages:
            user_id = msg.get("user", "unknown")
            text = msg.get("text", "")
            user_name = await get_user_name(user_id)
            attached_names = [
                str(file_data.get("name") or "attachment")
                for file_data in msg.get("files", [])
                if file_data.get("id")
            ]
            attachment_note = (
                f" [Attached: {', '.join(attached_names)}]"
                if attached_names
                else ""
            )
            thread_lines.append(f"{user_name}: {text}{attachment_note}")

        # Top-level message as title
        first_message = messages[0]
        message_title = first_message.get("text", "").strip()[:100] or "Slack thread"
        title = (
            f"#{channel_name} — {message_title}"
            if channel_name
            else message_title
        )

        body = "\n\n".join(thread_lines)

        # Permalink
        permalink = f"https://slack.com/archives/{channel}/p{thread_ts.replace('.', '')}"

        # Reference
        reference = f"slack:{channel}/{thread_ts}"

        # Source date
        source_date = datetime.utcfromtimestamp(float(thread_ts))

        items = [SourceItem(
            org_id=org_id,
            source=SourceType.SLACK,
            kind="thread",
            title=title,
            body=body,
            author=await get_user_name(first_message.get("user", "")),
            url=permalink,
            reference=reference,
            source_date=source_date,
            department_id=department_id,
        )]

        seen_files: set[str] = set()
        for message in messages:
            for file_data in message.get("files", []):
                file_id = str(file_data.get("id") or "")
                if not file_id or file_id in seen_files:
                    continue
                seen_files.add(file_id)
                filename = str(file_data.get("name") or f"{file_id}.txt")
                suffix = os.path.splitext(filename)[1].lower()
                if suffix not in SLACK_FILE_EXTENSIONS:
                    continue
                downloaded_name, content = await _download_slack_file(token, file_data)
                text = _extract_slack_file_text(content, downloaded_name)
                uploader_id = str(file_data.get("user") or message.get("user") or "")
                created = file_data.get("timestamp") or file_data.get("created")
                try:
                    file_date = datetime.utcfromtimestamp(float(created))
                except (TypeError, ValueError, OSError):
                    file_date = source_date
                items.append(SourceItem(
                    org_id=org_id,
                    source=SourceType.SLACK,
                    kind="file",
                    title=downloaded_name,
                    body=text,
                    author=await get_user_name(uploader_id),
                    url=str(file_data.get("permalink") or permalink),
                    reference=f"slack-file:{file_id}",
                    source_date=file_date,
                    department_id=department_id,
                ))

        return items


async def process_slack_events():
    """
    Worker: process unprocessed Slack events.

    Assembles threads and routes to extraction pipeline.
    """
    async with async_session() as session:
        from sqlalchemy import select

        # Get unprocessed message events
        result = await session.execute(
            select(EventRaw)
            .where(EventRaw.source == "slack")
            .where(EventRaw.event_type == "message")
            .where(EventRaw.processed_at.is_(None))
            .limit(100)
        )
        events = result.scalars().all()

        # Group by organization and thread because each organization owns a
        # distinct encrypted Slack token and channel allowlist.
        threads = defaultdict(list)

        for event in events:
            payload = event.payload
            event_data = payload.get("event", {})

            channel = event_data.get("channel")
            thread_ts = event_data.get("thread_ts") or event_data.get("ts")

            threads[(event.org_id, channel, thread_ts)].append(event)

        # Check which threads are ready (no messages in last 30 minutes)
        current_time = datetime.utcnow()

        for (org_id, channel, thread_ts), thread_events in threads.items():
            # Get latest message timestamp
            latest_ts = max(
                float(e.payload.get("event", {}).get("ts", 0))
                for e in thread_events
            )
            latest_dt = datetime.utcfromtimestamp(latest_ts)

            # Wait 30 minutes for thread to be quiet
            if (current_time - latest_dt) < timedelta(minutes=30):
                continue

            # Thread is ready - assemble and extract
            try:
                from persistence import list_connected_sources
                from pipelines.extract import extract_from_source

                sources = await list_connected_sources(org_id, include_config=True)
                source = next(
                    (item for item in sources if item["type"] == "slack"),
                    None,
                )
                config = source.get("config", {}) if source else {}
                watched_channels = config.get("watched_channels") or WATCHED_CHANNELS
                if not source or channel not in watched_channels:
                    for event in thread_events:
                        event.processed_at = datetime.utcnow()
                    await session.commit()
                    continue

                channel_names = config.get("channel_names") or {}
                source_items = await assemble_thread_items(
                    channel,
                    thread_ts,
                    org_id,
                    access_token=config.get("token"),
                    channel_name=channel_names.get(channel),
                    department_id=source.get("departmentId"),
                )

                for source_item in source_items:
                    await extract_from_source(source_item)
                    print(f"[Slack] Assembled source: {source_item.title}")

                # Mark all events in thread as processed
                for event in thread_events:
                    event.processed_at = datetime.utcnow()
                await session.commit()

            except Exception as e:
                for event in thread_events:
                    event.error = str(e)
                    event.processed_at = datetime.utcnow()
                await session.commit()
                print(f"[Slack] Error assembling thread: {e}")


async def fetch_slack_threads(
    org_id: str,
    *,
    access_token: str,
    channel_ids: List[str],
    channel_names: Optional[Dict[str, str]] = None,
    department_id: Optional[str] = None,
    days: int = 90,
    max_threads: int = 500,
) -> List[SourceItem]:
    """Fetch readable top-level messages and assemble their complete threads."""
    if not channel_ids:
        raise SlackApiError("Choose at least one Slack channel before syncing")

    oldest_ts = (datetime.utcnow() - timedelta(days=days)).timestamp()
    names = channel_names or {}
    source_items: list[SourceItem] = []
    seen_threads: set[tuple[str, str]] = set()
    assembled_threads = 0

    async with httpx.AsyncClient(
        headers=_headers(access_token),
        timeout=30,
    ) as client:
        for channel in channel_ids:
            cursor: Optional[str] = None
            while assembled_threads < max_threads:
                params: dict[str, Any] = {
                    "channel": channel,
                    "oldest": str(oldest_ts),
                    "limit": 100,
                }
                if cursor:
                    params["cursor"] = cursor
                payload = await _get(client, "conversations.history", params)

                for message in payload.get("messages", []):
                    thread_ts = str(message.get("thread_ts") or message.get("ts") or "")
                    identity = (channel, thread_ts)
                    if not thread_ts or identity in seen_threads:
                        continue
                    seen_threads.add(identity)
                    assembled_threads += 1
                    thread_items = await assemble_thread_items(
                        channel,
                        thread_ts,
                        org_id,
                        access_token=access_token,
                        channel_name=names.get(channel),
                        department_id=department_id,
                    )
                    source_items.extend(
                        item for item in thread_items if item.body.strip()
                    )
                    if assembled_threads >= max_threads:
                        break

                cursor = (
                    payload.get("response_metadata", {}).get("next_cursor")
                    or None
                )
                if not cursor:
                    break

    return sorted(
        source_items,
        key=lambda item: item.source_date or datetime.min,
        reverse=True,
    )


async def backfill_slack(
    org_id: str,
    days: int = 90,
    *,
    access_token: Optional[str] = None,
    channel_ids: Optional[List[str]] = None,
):
    """
    Backfill Slack history.

    Fetches messages from watched channels for the last N days.

    Args:
        org_id: Organization ID
        days: Days to backfill
    """
    source_items = await fetch_slack_threads(
        org_id,
        access_token=access_token or SLACK_BOT_TOKEN,
        channel_ids=channel_ids or WATCHED_CHANNELS,
        days=days,
    )
    print(f"Slack backfill complete for org {org_id}: {len(source_items)} threads")
    return source_items


if __name__ == "__main__":
    import asyncio

    async def main():
        await backfill_slack(
            org_id="test-org",
            days=30
        )

    asyncio.run(main())
